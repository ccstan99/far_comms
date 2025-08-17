#!/usr/bin/env python

import logging
import glob
import os
import base64
from pathlib import Path
from langchain_community.document_loaders import AssemblyAIAudioTranscriptLoader, PyPDFLoader
from langchain_community.document_loaders.assemblyai import TranscriptFormat

logger = logging.getLogger(__name__)


def _decode_qr_codes_from_image(img_data: bytes) -> list:
    """
    Decode QR codes from image data using pyzbar.
    
    Args:
        img_data: PNG image bytes
        
    Returns:
        List of QR code URLs found
    """
    try:
        from pyzbar import pyzbar
        from PIL import Image
        import io
        
        # Convert bytes to PIL Image
        image = Image.open(io.BytesIO(img_data))
        
        # Decode QR codes
        qr_codes = []
        decoded_objects = pyzbar.decode(image)
        
        for obj in decoded_objects:
            if obj.type == 'QRCODE':
                qr_data = obj.data.decode('utf-8')
                qr_codes.append(qr_data)
                
        return qr_codes
        
    except ImportError:
        logger.warning("pyzbar not available - cannot decode QR codes")
        return []
    except Exception as e:
        logger.warning(f"QR code decoding failed: {e}")
        return []


def _analyze_pdf_visually(pdf_path: str, speaker_name: str = None) -> dict:
    """
    Analyze PDF visually using multimodal LLM to extract QR codes and describe images.
    Also saves image-rich slides for potential social media use.
    
    Args:
        pdf_path: Path to PDF file
        speaker_name: Speaker name for image filename generation
        
    Returns:
        dict with qr_codes, images, visual_elements, and saved_images
    """
    try:
        import fitz  # PyMuPDF
        from anthropic import Anthropic
        import os
        
        # Initialize Anthropic client
        api_key = os.getenv("ANTHROPIC_API_KEY")
        if not api_key:
            # Try loading from .env file
            try:
                from dotenv import load_dotenv
                load_dotenv()
                api_key = os.getenv("ANTHROPIC_API_KEY")
            except ImportError:
                pass
                
        if not api_key:
            logger.warning("ANTHROPIC_API_KEY not found - skipping visual analysis")
            return {"qr_codes": [], "visual_elements": [], "page_analyses": []}
        client = Anthropic(api_key=api_key)
        
        doc = fitz.open(pdf_path)
        results = {
            "qr_codes": [],
            "visual_elements": [],
            "page_analyses": [],
            "saved_images": []
        }
        
        # Set up output directory for images
        from far_comms.utils.project_paths import get_output_dir
        output_dir = get_output_dir()
        
        # Clean speaker name for filename
        safe_speaker_name = "unknown"
        if speaker_name:
            import re
            safe_speaker_name = re.sub(r'[^\w\-_\s]', '', speaker_name.replace(' ', '_'))
            safe_speaker_name = re.sub(r'[_\s]+', '_', safe_speaker_name).strip('_')
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            
            # Convert page to image
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x zoom for better quality
            img_data = pix.tobytes("png")
            img_base64 = base64.b64encode(img_data).decode()
            
            # Save full slide image for easy access
            slide_filename = f"slide_{page_num + 1:02d}.png"
            slide_path = output_dir / slide_filename
            
            try:
                with open(slide_path, "wb") as slide_file:
                    slide_file.write(img_data)
                logger.info(f"Saved full slide {page_num + 1} as {slide_filename}")
            except Exception as save_error:
                logger.warning(f"Failed to save full slide {page_num + 1}: {save_error}")
            
            # First try to decode QR codes directly using pyzbar
            qr_urls = _decode_qr_codes_from_image(img_data)
            
            # Analyze with multimodal LLM
            prompt = """Analyze this slide image and extract:
1. QR codes: If you see any QR codes, try to read the URL they contain
2. Visual elements: Describe any charts, diagrams, tables, or images with brief alt text
3. Important text: Any key text that might be missed by OCR
4. Image richness: BE VERY SELECTIVE - only for slides with compelling visual data

STRICT CRITERIA FOR "is_image_rich": Only mark as "true" if slide contains:
✅ Complex workflow diagrams with arrows/boxes/connections (like process flows)
✅ Data tables with numbers/results/metrics (not just text lists)  
✅ Charts/graphs with data visualization or performance comparisons
✅ Technical system diagrams with visual components
✅ Comparison tables showing quantitative results

❌ DO NOT mark as image-rich:
❌ Title slides with just names/affiliations/logos
❌ Bullet point lists (even with fancy formatting)
❌ Text-heavy slides with minimal visuals
❌ Simple layouts that are mostly text
❌ Basic organizational info or contact details

Format your response as JSON:
{
  "qr_codes": [{"url": "detected_url", "location": "description of where on slide"}],
  "visual_elements": [{"type": "chart|diagram|table|image", "description": "brief alt text"}],
  "key_text": ["any important text visible"],
  "is_image_rich": "true|false - ONLY true for slides with quantitative data, complex diagrams, or rich visual content worth sharing on social media",
  "social_media_potential": "brief explanation focusing on visual complexity and data presentation value"
}"""

            try:
                response = client.messages.create(
                    model="claude-sonnet-4-20250514",
                    max_tokens=1024,
                    messages=[{
                        "role": "user",
                        "content": [
                            {"type": "text", "text": prompt},
                            {"type": "image", "source": {
                                "type": "base64",
                                "media_type": "image/png",
                                "data": img_base64
                            }}
                        ]
                    }]
                )
                
                # Parse response
                import json
                analysis_text = response.content[0].text
                
                # Extract JSON from response
                if "{" in analysis_text and "}" in analysis_text:
                    json_start = analysis_text.find("{")
                    json_end = analysis_text.rfind("}") + 1
                    json_str = analysis_text[json_start:json_end]
                    
                    analysis = json.loads(json_str)
                    
                    # Add decoded QR codes (real URLs from pyzbar)
                    for qr_url in qr_urls:
                        results["qr_codes"].append({
                            "url": qr_url,
                            "location": f"QR code on page {page_num + 1}",
                            "page": page_num + 1,
                            "source": "pyzbar_decoded"
                        })
                    
                    # Add any QR codes detected by Claude (for location info)
                    for qr in analysis.get("qr_codes", []):
                        if not qr_urls:  # Only add if pyzbar didn't find any
                            qr["page"] = page_num + 1
                            qr["source"] = "claude_detected"
                            results["qr_codes"].append(qr)
                    
                    for element in analysis.get("visual_elements", []):
                        element["page"] = page_num + 1
                        results["visual_elements"].append(element)
                    
                    # Save image if it's rich in visual content
                    is_image_rich = analysis.get("is_image_rich", "false").lower() == "true"
                    if is_image_rich:
                        try:
                            # Save the slide image
                            image_filename = f"{safe_speaker_name}_{page_num + 1}.png"
                            image_path = output_dir / image_filename
                            
                            with open(image_path, "wb") as img_file:
                                img_file.write(img_data)
                            
                            results["saved_images"].append({
                                "page": page_num + 1,
                                "filename": image_filename,
                                "path": str(image_path),
                                "social_media_potential": analysis.get("social_media_potential", ""),
                                "visual_elements_count": len(analysis.get("visual_elements", []))
                            })
                            
                            logger.info(f"Saved image-rich slide {page_num + 1} as {image_filename}")
                            
                        except Exception as save_error:
                            logger.warning(f"Failed to save slide {page_num + 1} image: {save_error}")
                    
                    results["page_analyses"].append({
                        "page": page_num + 1,
                        "analysis": analysis
                    })
                
            except Exception as e:
                logger.warning(f"Failed to analyze page {page_num + 1} of PDF: {e}")
                # Still add QR codes even if visual analysis fails
                for qr_url in qr_urls:
                    results["qr_codes"].append({
                        "url": qr_url,
                        "location": f"QR code on page {page_num + 1}",
                        "page": page_num + 1,
                        "source": "pyzbar_decoded"
                    })
                # Add a note about failed analysis
                results["page_analyses"].append({
                    "page": page_num + 1,
                    "analysis": {"error": f"Visual analysis failed: {str(e)}", "qr_codes": qr_urls}
                })
                continue
        
        doc.close()
        
        logger.info(f"Visual analysis complete: {len(results['qr_codes'])} QR codes, {len(results['visual_elements'])} visual elements, {len(results['saved_images'])} images saved")
        return results
        
    except ImportError:
        logger.warning("PyMuPDF not available - skipping visual analysis")
        return {"qr_codes": [], "visual_elements": [], "page_analyses": [], "saved_images": []}
    except Exception as e:
        logger.error(f"Error in visual analysis: {e}")
        return {"qr_codes": [], "visual_elements": [], "page_analyses": [], "saved_images": []}


def convert_pptx_to_pdf(pptx_path: str, output_dir: str = None) -> str:
    """
    Convert PowerPoint (PPTX/PPT) to PDF using python-pptx and reportlab.
    
    Args:
        pptx_path: Path to the PowerPoint file
        output_dir: Directory to save the PDF (defaults to same dir as PPTX)
    
    Returns:
        str: Path to the generated PDF file
        
    Raises:
        Exception: If conversion fails
    """
    try:
        from pptx import Presentation
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter, A4
        from reportlab.lib.units import inch
        from reportlab.lib.utils import ImageReader
        from io import BytesIO
        import tempfile
        
        logger.info(f"Converting PPTX to PDF: {pptx_path}")
        
        # Load PowerPoint presentation
        prs = Presentation(pptx_path)
        
        # Determine output path
        if output_dir is None:
            output_dir = os.path.dirname(pptx_path)
        
        output_filename = Path(pptx_path).stem + "_converted.pdf"
        output_path = os.path.join(output_dir, output_filename)
        
        # Create PDF canvas
        c = canvas.Canvas(output_path, pagesize=A4)
        width, height = A4
        
        for slide_idx, slide in enumerate(prs.slides):
            logger.info(f"Processing slide {slide_idx + 1}/{len(prs.slides)}")
            
            # Add slide title if present
            y_position = height - 72  # Start 1 inch from top
            
            # Extract text from slide
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            
            # Add slide number
            c.setFont("Helvetica-Bold", 16)
            c.drawString(72, y_position, f"Slide {slide_idx + 1}")
            y_position -= 36
            
            # Add text content
            c.setFont("Helvetica", 12)
            for text in slide_texts:
                if y_position < 72:  # If near bottom of page, start new page
                    c.showPage()
                    y_position = height - 72
                
                # Split long text into lines
                lines = text.split('\n')
                for line in lines:
                    if len(line) > 80:  # Wrap long lines
                        words = line.split(' ')
                        current_line = ""
                        for word in words:
                            if len(current_line + " " + word) > 80:
                                if current_line:
                                    c.drawString(72, y_position, current_line)
                                    y_position -= 18
                                current_line = word
                            else:
                                current_line += (" " + word if current_line else word)
                        if current_line:
                            c.drawString(72, y_position, current_line)
                            y_position -= 18
                    else:
                        c.drawString(72, y_position, line)
                        y_position -= 18
                
                y_position -= 12  # Extra space between shapes
            
            # Start new page for next slide (except for last slide)
            if slide_idx < len(prs.slides) - 1:
                c.showPage()
        
        # Save PDF
        c.save()
        logger.info(f"Successfully converted PPTX to PDF: {output_path}")
        return output_path
        
    except Exception as e:
        logger.error(f"Failed to convert PPTX to PDF: {e}")
        raise Exception(f"PPTX conversion failed: {str(e)}")


def find_presentation(speaker_name: str) -> tuple[str, str]:
    """Find PDF or PPTX file that best matches speaker name
    
    Returns:
        tuple: (file_path, file_type) where file_type is 'pdf' or 'pptx'
    """
    # Search for both PDF and PPTX files
    pdf_files = glob.glob("data/slides/*.pdf")
    pptx_files = glob.glob("data/slides/*.pptx") + glob.glob("data/slides/*.ppt")
    
    all_files = [(f, 'pdf') for f in pdf_files] + [(f, 'pptx') for f in pptx_files]
    
    if not all_files:
        return None, None
        
    # Split speaker name into individual words
    speaker_words = [word.lower() for word in speaker_name.replace("-", " ").replace("_", " ").split() if len(word) > 2]
    
    best_match = None
    best_type = None
    best_score = 0
    
    for file_path, file_type in all_files:
        filename = os.path.basename(file_path).lower()
        score = 0
        
        # Score based on how many speaker name words appear in filename
        for word in speaker_words:
            if word in filename:
                score += len(word)
        
        # Prefer PDF files by adding a small bonus (they process more reliably)
        if file_type == 'pdf':
            score += 0.5
        
        if score > best_score:
            best_score = score
            best_match = file_path
            best_type = file_type
    
    return (best_match, best_type) if best_score > 0 else (None, None)


def find_pdf(speaker_name: str) -> str:
    """Legacy function - find PDF file that best matches speaker name"""
    file_path, file_type = find_presentation(speaker_name)
    return file_path if file_type == 'pdf' else None


def extract_pdf(pdf_path: str, speaker_name: str = None) -> dict:
    """Extract both text and visual content from PDF"""
    try:
        # Extract text content
        text_loader = PyPDFLoader(pdf_path)
        text_docs = text_loader.load()
        text_content = "\n\n".join([doc.page_content for doc in text_docs])
        
        logger.info(f"PyPDFLoader extracted {len(text_docs)} pages with {len(text_content)} chars")
        
        # Extract visual content (QR codes, images, charts) and save image-rich slides
        visual_analysis = _analyze_pdf_visually(pdf_path, speaker_name)
        
        # Check for missing pages by comparing PDF page count with extracted pages
        import fitz
        doc = fitz.open(pdf_path)
        total_pdf_pages = len(doc)
        doc.close()
        
        if len(text_docs) != total_pdf_pages:
            logger.warning(f"Page count mismatch: PDF has {total_pdf_pages} pages but PyPDFLoader extracted {len(text_docs)} pages")
        
        # Combine text with visual descriptions
        enhanced_content = text_content
        if visual_analysis["visual_elements"]:
            enhanced_content += "\n\n--- VISUAL ELEMENTS ---\n"
            for element in visual_analysis["visual_elements"]:
                enhanced_content += f"[Page {element['page']}] {element['type']}: {element['description']}\n"
        
        return {
            "text_content": text_content,
            "enhanced_content": enhanced_content,
            "visual_analysis": visual_analysis,
            "qr_codes": visual_analysis["qr_codes"],
            "visual_elements": visual_analysis["visual_elements"],
            "saved_images": visual_analysis["saved_images"],
            "page_count_info": f"PDF: {total_pdf_pages} pages, Extracted: {len(text_docs)} pages"
        }
    except Exception as e:
        logger.error(f"Error extracting PDF content from {pdf_path}: {e}")
        return {
            "text_content": "",
            "enhanced_content": "",
            "visual_analysis": {"qr_codes": [], "visual_elements": [], "page_analyses": [], "saved_images": []},
            "qr_codes": [],
            "visual_elements": [],
            "saved_images": []
        }


def find_video(speaker_name: str) -> str:
    """Find video file that best matches speaker name"""
    video_patterns = ["*.mp4", "*.mkv", "*.avi", "*.mov", "*.webm", "*.m4a", "*.wav"]
    video_files = []
    
    for pattern in video_patterns:
        video_files.extend(glob.glob(f"data/videos/{pattern}"))
    
    if not video_files:
        return None
        
    # Split speaker name into individual words
    speaker_words = [word.lower() for word in speaker_name.replace("-", " ").replace("_", " ").split() if len(word) > 2]
    
    best_match = None
    best_score = 0
    
    for video_path in video_files:
        filename = os.path.basename(video_path).lower()
        score = 0
        
        # Score based on how many speaker name words appear in filename
        for word in speaker_words:
            if word in filename:
                score += len(word)
        
        if score > best_score:
            best_score = score
            best_match = video_path
    
    return best_match if best_score > 0 else None


def extract_youtube(youtube_url: str) -> dict:
    """Extract transcript from YouTube using AssemblyAI with yt-dlp fallback"""
    try:
        # First try direct URL with AssemblyAI
        loader = AssemblyAIAudioTranscriptLoader(
            file_path=youtube_url,
            transcript_format=TranscriptFormat.SUBTITLES_SRT
        )
        docs = loader.load()
        
        if docs:
            srt_content = docs[0].page_content
            return {
                "success": True,
                "srt_content": srt_content,
                "source": "youtube_url_direct"
            }
        else:
            return {"success": False, "error": "No content returned from AssemblyAI"}
            
    except Exception as e:
        error_str = str(e)
        logger.warning(f"Direct YouTube URL failed: {error_str}")
        
        # If we get HTML/text error, try yt-dlp download first
        if "text/html" in error_str or "HTML document" in error_str:
            try:
                import yt_dlp
                import tempfile
                import os
                
                with tempfile.TemporaryDirectory() as temp_dir:
                    audio_path = os.path.join(temp_dir, "audio.%(ext)s")
                    
                    ydl_opts = {
                        'format': 'bestaudio[ext=m4a]/bestaudio[ext=mp3]/bestaudio/best[height<=480]',
                        'outtmpl': audio_path,
                        'extract_flat': False,
                        'ignoreerrors': True,
                        'no_warnings': False,
                        'extractaudio': True,
                        'audioformat': 'mp3',
                        'embed_subs': False,
                        'writesubtitles': False,
                        'writeautomaticsub': False,
                        # Anti-bot measures
                        'user_agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36',
                        'referer': 'https://www.youtube.com/',
                        'http_chunk_size': 10485760,
                        'retries': 3,
                        'fragment_retries': 3,
                        'skip_unavailable_fragments': True,
                    }
                    
                    with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                        ydl.download([youtube_url])
                    
                    # Find the downloaded file
                    audio_files = [f for f in os.listdir(temp_dir) if f.startswith("audio.")]
                    if not audio_files:
                        return {"success": False, "error": "yt-dlp download failed - no audio file created"}
                    
                    downloaded_file = os.path.join(temp_dir, audio_files[0])
                    
                    # Now transcribe the downloaded file
                    loader = AssemblyAIAudioTranscriptLoader(
                        file_path=downloaded_file,
                        transcript_format=TranscriptFormat.SUBTITLES_SRT
                    )
                    docs = loader.load()
                    
                    if docs:
                        return {
                            "success": True,
                            "srt_content": docs[0].page_content,
                            "source": "youtube_url_yt-dlp"
                        }
                    else:
                        return {"success": False, "error": "AssemblyAI transcription failed on downloaded file"}
                        
            except ImportError:
                logger.error("yt-dlp not available for YouTube download fallback")
                return {"success": False, "error": f"YouTube URL failed: {error_str} (yt-dlp not available)"}
            except Exception as fallback_error:
                logger.error(f"yt-dlp fallback failed: {fallback_error}")
                return {"success": False, "error": f"Both direct and yt-dlp failed: {error_str}, {str(fallback_error)}"}
        else:
            return {"success": False, "error": error_str}


def extract_video(video_path: str) -> dict:
    """Extract transcript from local video file using AssemblyAI"""
    try:
        # Use SRT format for timestamps
        loader = AssemblyAIAudioTranscriptLoader(
            file_path=video_path,
            transcript_format=TranscriptFormat.SUBTITLES_SRT
        )
        docs = loader.load()
        
        if docs:
            srt_content = docs[0].page_content
            return {
                "success": True,
                "srt_content": srt_content,
                "source": "local_video"
            }
        else:
            return {"success": False, "error": "No content returned from AssemblyAI"}
            
    except Exception as e:
        logger.error(f"Error extracting local video transcript: {e}")
        return {"success": False, "error": str(e)}