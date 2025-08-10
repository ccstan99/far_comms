#!/usr/bin/env python

import os
import base64
from pathlib import Path
from typing import Optional, Dict, Any, List
import logging

try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False

try:
    from pptx import Presentation
    PYTHON_PPTX_AVAILABLE = True
except ImportError:
    PYTHON_PPTX_AVAILABLE = False

logger = logging.getLogger(__name__)


def get_slide_images(file_path: str, max_slides: int = 10) -> Dict[str, Any]:
    """
    Extract slide images from PDF for visual analysis
    
    Args:
        file_path: Path to the PDF file
        max_slides: Maximum number of slides to extract (cost control)
    
    Returns:
        Dictionary containing slide images and metadata
    """
    if not PYMUPDF_AVAILABLE:
        return {
            "error": "PyMuPDF not available. Install with: pip install PyMuPDF",
            "images": [],
            "page_count": 0,
            "success": False
        }
    
    try:
        doc = fitz.open(file_path)
        images = []
        page_count = len(doc)
        
        # Limit slides for cost control (unless max_slides is negative, meaning process all)
        if max_slides < 0:
            slides_to_process = page_count  # Process all slides
        else:
            slides_to_process = min(max_slides, page_count)
        
        for page_num in range(slides_to_process):
            page = doc.load_page(page_num)
            
            # Convert page to image (2x scale for better quality)
            mat = fitz.Matrix(2, 2)  # 2x zoom
            pix = page.get_pixmap(matrix=mat)
            img_data = pix.tobytes("png")
            
            # Convert to base64 for LLM processing
            img_base64 = base64.b64encode(img_data).decode('utf-8')
            
            images.append({
                "page": page_num + 1,
                "image_base64": img_base64,
                "width": pix.width,
                "height": pix.height
            })
            
            logger.debug(f"Extracted image for slide {page_num + 1}: {pix.width}x{pix.height}")
        
        doc.close()
        
        return {
            "file_type": "pdf",
            "file_name": Path(file_path).name,
            "page_count": page_count,
            "slides_processed": slides_to_process,
            "images": images,
            "success": True
        }
        
    except Exception as e:
        logger.error(f"Error extracting slide images from {file_path}: {e}")
        return {
            "error": str(e),
            "images": [],
            "page_count": 0,
            "success": False
        }


def get_pdf_text(file_path: str) -> Dict[str, Any]:
    """
    Extract text content from a PDF file using PyMuPDF
    
    Args:
        file_path: Path to the PDF file
    
    Returns:
        Dictionary containing extracted content and metadata
    """
    if not PYMUPDF_AVAILABLE:
        return {
            "error": "PyMuPDF not available. Install with: pip install PyMuPDF",
            "content": "",
            "slides": [],
            "page_count": 0
        }
    
    try:
        doc = fitz.open(file_path)
        slides = []
        all_text = []
        page_count = len(doc)  # Store this before closing
        
        for page_num in range(page_count):
            page = doc.load_page(page_num)
            text = page.get_text()
            
            # Always include pages to maintain structure
            slides.append({
                "page": page_num + 1,
                "content": text.strip() if text.strip() else f"[PAGE {page_num + 1} - NO TEXT EXTRACTED]"
            })
            
            if text.strip():
                all_text.append(text.strip())
        
        doc.close()
        
        full_content = "\n\n---PAGE BREAK---\n\n".join(all_text)
        
        return {
            "file_type": "pdf",
            "file_name": Path(file_path).name,
            "page_count": page_count,
            "slides": slides,
            "content": full_content,
            "success": True
        }
        
    except Exception as e:
        logger.error(f"Error extracting PDF content from {file_path}: {e}")
        return {
            "error": str(e),
            "content": "",
            "slides": [],
            "page_count": 0,
            "success": False
        }


def get_pptx_text(file_path: str) -> Dict[str, Any]:
    """
    Extract text content from a PowerPoint file using python-pptx
    
    Args:
        file_path: Path to the PPTX file
    
    Returns:
        Dictionary containing extracted content and metadata
    """
    if not PYTHON_PPTX_AVAILABLE:
        return {
            "error": "python-pptx not available. Install with: pip install python-pptx",
            "content": "",
            "slides": [],
            "slide_count": 0
        }
    
    try:
        prs = Presentation(file_path)
        slides = []
        all_text = []
        
        for slide_num, slide in enumerate(prs.slides):
            slide_text = []
            
            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            if slide_text:  # Only include slides with content
                slide_content = "\n".join(slide_text)
                slides.append({
                    "slide": slide_num + 1,
                    "content": slide_content
                })
                all_text.append(slide_content)
        
        full_content = "\n\n---SLIDE BREAK---\n\n".join(all_text)
        
        return {
            "file_type": "pptx",
            "file_name": Path(file_path).name,
            "slide_count": len(prs.slides),
            "slides": slides,
            "content": full_content,
            "success": True
        }
        
    except Exception as e:
        logger.error(f"Error extracting PPTX content from {file_path}: {e}")
        return {
            "error": str(e),
            "content": "",
            "slides": [],
            "slide_count": 0,
            "success": False
        }


def get_slide_text(file_path: str) -> Dict[str, Any]:
    """
    Extract content from presentation files (PDF or PPTX)
    
    Args:
        file_path: Path to the presentation file
    
    Returns:
        Dictionary containing extracted content and metadata
    """
    if not os.path.exists(file_path):
        return {
            "error": f"File not found: {file_path}",
            "content": "",
            "slides": [],
            "success": False
        }
    
    file_path = str(Path(file_path).resolve())
    file_ext = Path(file_path).suffix.lower()
    
    if file_ext == '.pdf':
        return get_pdf_text(file_path)
    elif file_ext in ['.pptx', '.ppt']:
        return get_pptx_text(file_path)
    else:
        return {
            "error": f"Unsupported file type: {file_ext}. Supported: .pdf, .pptx, .ppt",
            "content": "",
            "slides": [],
            "success": False
        }


def get_slide_content(file_path: str, max_slides: int = 10) -> Dict[str, Any]:
    """
    Extract both text content and images from presentation files
    
    Args:
        file_path: Path to the presentation file
        max_slides: Maximum number of slides to extract images for (cost control)
    
    Returns:
        Dictionary containing both text and image data
    """
    if not os.path.exists(file_path):
        return {
            "error": f"File not found: {file_path}",
            "content": "",
            "slides": [],
            "images": [],
            "success": False
        }
    
    file_path = str(Path(file_path).resolve())
    file_ext = Path(file_path).suffix.lower()
    
    if file_ext == '.pdf':
        # Extract text content
        text_result = get_pdf_text(file_path)
        
        # Extract images
        image_result = get_slide_images(file_path, max_slides)
        
        if text_result.get("success") and image_result.get("success"):
            # Combine both results
            combined_result = text_result.copy()
            combined_result.update({
                "images": image_result.get("images", []),
                "slides_processed_for_images": image_result.get("slides_processed", 0),
                "has_visual_data": True
            })
            return combined_result
        elif text_result.get("success"):
            # Text extraction succeeded, image extraction failed
            text_result.update({
                "images": [],
                "has_visual_data": False,
                "image_extraction_error": image_result.get("error", "Unknown error")
            })
            return text_result
        else:
            # Both failed
            return text_result
    
    elif file_ext in ['.pptx', '.ppt']:
        # For now, only text extraction for PPTX (images could be added later)
        result = get_pptx_text(file_path)
        result.update({
            "images": [],
            "has_visual_data": False,
            "image_extraction_note": "Image extraction not yet implemented for PPTX files"
        })
        return result
    
    else:
        return {
            "error": f"Unsupported file type: {file_ext}. Supported: .pdf, .pptx, .ppt",
            "content": "",
            "slides": [],
            "images": [],
            "success": False
        }


def format_extraction_summary(content_dict: Dict[str, Any]) -> str:
    """
    Get a concise summary of extracted slide content
    
    Args:
        content_dict: Result from get_slide_text()
    
    Returns:
        Human-readable summary string
    """
    if not content_dict.get("success"):
        return f"Extraction failed: {content_dict.get('error', 'Unknown error')}"
    
    file_type = content_dict.get("file_type", "unknown")
    file_name = content_dict.get("file_name", "unknown")
    
    if file_type == "pdf":
        count = content_dict.get("page_count", 0)
        unit = "pages"
    else:
        count = content_dict.get("slide_count", 0)
        unit = "slides"
    
    content_length = len(content_dict.get("content", ""))
    
    return f"Extracted from {file_name}: {count} {unit}, {content_length} characters"


# Example usage and testing
if __name__ == "__main__":
    # Test with the provided file
    test_file = "/Users/cheng2/Desktop/agents/far_comms/data/slides/11_50_Xiaoyuan Yi-ValueCompass_updated.pptx.pdf"
    
    print("Testing slide extraction...")
    result = get_slide_text(test_file)
    
    print(f"Success: {result.get('success')}")
    print(format_extraction_summary(result))
    
    if result.get("success"):
        print(f"\nFirst 500 characters of content:")
        print(result.get("content", "")[:500])
        print("\n" + "="*50)
        
        print(f"\nFirst slide/page content:")
        slides = result.get("slides", [])
        if slides:
            print(slides[0].get("content", "")[:300])
    else:
        print(f"Error: {result.get('error')}")