#!/usr/bin/env python

import os
import re
import base64
from pathlib import Path
from typing import Optional, Dict, Any, List
import logging
from crewai import LLM

try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False

try:
    from pptx import Presentation
    PYTHON_PPTX_AVAILABLE = True
except ImportError:
    PYTHON_PPTX_AVAILABLE = False

logger = logging.getLogger(__name__)


def get_slide_images(file_path: str, max_slides: int = 10) -> Dict[str, Any]:
    """
    Extract slide images from PDF for visual analysis
    
    Args:
        file_path: Path to the PDF file
        max_slides: Maximum number of slides to extract (cost control)
    
    Returns:
        Dictionary containing slide images and metadata
    """
    if not PYMUPDF_AVAILABLE:
        return {
            "error": "PyMuPDF not available. Install with: pip install PyMuPDF",
            "images": [],
            "page_count": 0,
            "success": False
        }
    
    try:
        doc = fitz.open(file_path)
        images = []
        page_count = len(doc)
        
        # Limit slides for cost control (unless max_slides is negative, meaning process all)
        if max_slides < 0:
            slides_to_process = page_count  # Process all slides
        else:
            slides_to_process = min(max_slides, page_count)
        
        for page_num in range(slides_to_process):
            page = doc.load_page(page_num)
            
            # Convert page to image (2x scale for better quality)
            mat = fitz.Matrix(2, 2)  # 2x zoom
            pix = page.get_pixmap(matrix=mat)
            img_data = pix.tobytes("png")
            
            # Convert to base64 for LLM processing
            img_base64 = base64.b64encode(img_data).decode('utf-8')
            
            images.append({
                "page": page_num + 1,
                "image_base64": img_base64,
                "width": pix.width,
                "height": pix.height
            })
            
            logger.debug(f"Extracted image for slide {page_num + 1}: {pix.width}x{pix.height}")
        
        doc.close()
        
        return {
            "file_type": "pdf",
            "file_name": Path(file_path).name,
            "page_count": page_count,
            "slides_processed": slides_to_process,
            "images": images,
            "success": True
        }
        
    except Exception as e:
        logger.error(f"Error extracting slide images from {file_path}: {e}")
        return {
            "error": str(e),
            "images": [],
            "page_count": 0,
            "success": False
        }


def get_pdf_text(file_path: str) -> Dict[str, Any]:
    """
    Extract text content from a PDF file using PyMuPDF
    
    Args:
        file_path: Path to the PDF file
    
    Returns:
        Dictionary containing extracted content and metadata
    """
    if not PYMUPDF_AVAILABLE:
        return {
            "error": "PyMuPDF not available. Install with: pip install PyMuPDF",
            "content": "",
            "slides": [],
            "page_count": 0
        }
    
    try:
        doc = fitz.open(file_path)
        slides = []
        all_text = []
        page_count = len(doc)  # Store this before closing
        
        for page_num in range(page_count):
            page = doc.load_page(page_num)
            text = page.get_text()
            
            # Always include pages to maintain structure
            slides.append({
                "page": page_num + 1,
                "content": text.strip() if text.strip() else f"[PAGE {page_num + 1} - NO TEXT EXTRACTED]"
            })
            
            if text.strip():
                all_text.append(text.strip())
        
        doc.close()
        
        full_content = "\n\n---PAGE BREAK---\n\n".join(all_text)
        
        return {
            "file_type": "pdf",
            "file_name": Path(file_path).name,
            "page_count": page_count,
            "slides": slides,
            "content": full_content,
            "success": True
        }
        
    except Exception as e:
        logger.error(f"Error extracting PDF content from {file_path}: {e}")
        return {
            "error": str(e),
            "content": "",
            "slides": [],
            "page_count": 0,
            "success": False
        }


def get_pptx_text(file_path: str) -> Dict[str, Any]:
    """
    Extract text content from a PowerPoint file using python-pptx
    
    Args:
        file_path: Path to the PPTX file
    
    Returns:
        Dictionary containing extracted content and metadata
    """
    if not PYTHON_PPTX_AVAILABLE:
        return {
            "error": "python-pptx not available. Install with: pip install python-pptx",
            "content": "",
            "slides": [],
            "slide_count": 0
        }
    
    try:
        prs = Presentation(file_path)
        slides = []
        all_text = []
        
        for slide_num, slide in enumerate(prs.slides):
            slide_text = []
            
            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            if slide_text:  # Only include slides with content
                slide_content = "\n".join(slide_text)
                slides.append({
                    "slide": slide_num + 1,
                    "content": slide_content
                })
                all_text.append(slide_content)
        
        full_content = "\n\n---SLIDE BREAK---\n\n".join(all_text)
        
        return {
            "file_type": "pptx",
            "file_name": Path(file_path).name,
            "slide_count": len(prs.slides),
            "slides": slides,
            "content": full_content,
            "success": True
        }
        
    except Exception as e:
        logger.error(f"Error extracting PPTX content from {file_path}: {e}")
        return {
            "error": str(e),
            "content": "",
            "slides": [],
            "slide_count": 0,
            "success": False
        }


def get_slide_text(file_path: str) -> Dict[str, Any]:
    """
    Extract content from presentation files (PDF or PPTX)
    
    Args:
        file_path: Path to the presentation file
    
    Returns:
        Dictionary containing extracted content and metadata
    """
    if not os.path.exists(file_path):
        return {
            "error": f"File not found: {file_path}",
            "content": "",
            "slides": [],
            "success": False
        }
    
    file_path = str(Path(file_path).resolve())
    file_ext = Path(file_path).suffix.lower()
    
    if file_ext == '.pdf':
        return get_pdf_text(file_path)
    elif file_ext in ['.pptx', '.ppt']:
        return get_pptx_text(file_path)
    else:
        return {
            "error": f"Unsupported file type: {file_ext}. Supported: .pdf, .pptx, .ppt",
            "content": "",
            "slides": [],
            "success": False
        }


def get_slide_content(file_path: str, max_slides: int = 10) -> Dict[str, Any]:
    """
    Extract both text content and images from presentation files
    
    Args:
        file_path: Path to the presentation file
        max_slides: Maximum number of slides to extract images for (cost control)
    
    Returns:
        Dictionary containing both text and image data
    """
    if not os.path.exists(file_path):
        return {
            "error": f"File not found: {file_path}",
            "content": "",
            "slides": [],
            "images": [],
            "success": False
        }
    
    file_path = str(Path(file_path).resolve())
    file_ext = Path(file_path).suffix.lower()
    
    if file_ext == '.pdf':
        # Extract text content
        text_result = get_pdf_text(file_path)
        
        # Extract images
        image_result = get_slide_images(file_path, max_slides)
        
        if text_result.get("success") and image_result.get("success"):
            # Combine both results
            combined_result = text_result.copy()
            combined_result.update({
                "images": image_result.get("images", []),
                "slides_processed_for_images": image_result.get("slides_processed", 0),
                "has_visual_data": True
            })
            return combined_result
        elif text_result.get("success"):
            # Text extraction succeeded, image extraction failed
            text_result.update({
                "images": [],
                "has_visual_data": False,
                "image_extraction_error": image_result.get("error", "Unknown error")
            })
            return text_result
        else:
            # Both failed
            return text_result
    
    elif file_ext in ['.pptx', '.ppt']:
        # For now, only text extraction for PPTX (images could be added later)
        result = get_pptx_text(file_path)
        result.update({
            "images": [],
            "has_visual_data": False,
            "image_extraction_note": "Image extraction not yet implemented for PPTX files"
        })
        return result
    
    else:
        return {
            "error": f"Unsupported file type: {file_ext}. Supported: .pdf, .pptx, .ppt",
            "content": "",
            "slides": [],
            "images": [],
            "success": False
        }


def format_extraction_summary(content_dict: Dict[str, Any]) -> str:
    """
    Get a concise summary of extracted slide content
    
    Args:
        content_dict: Result from get_slide_text()
    
    Returns:
        Human-readable summary string
    """
    if not content_dict.get("success"):
        return f"Extraction failed: {content_dict.get('error', 'Unknown error')}"
    
    file_type = content_dict.get("file_type", "unknown")
    file_name = content_dict.get("file_name", "unknown")
    
    if file_type == "pdf":
        count = content_dict.get("page_count", 0)
        unit = "pages"
    else:
        count = content_dict.get("slide_count", 0)
        unit = "slides"
    
    content_length = len(content_dict.get("content", ""))
    
    return f"Extracted from {file_name}: {count} {unit}, {content_length} characters"


# Text cleaning functions (formerly in slide_cleaner.py)

def clean_slide_text(raw_content: str) -> str:
    """
    Use LLM to clean up messy slide text extraction artifacts
    
    Args:
        raw_content: Raw text extracted from PDF/PPTX slides
    
    Returns:
        Cleaned, readable text suitable for content generation
    """
    if not raw_content or not raw_content.strip():
        return raw_content
    
    # If content is too short, probably clean enough already
    if len(raw_content.strip()) < 100:
        return raw_content.strip()
    
    try:
        llm = LLM(
            model="anthropic/claude-haiku-3-20240307",  # Fast, cheap model for cleaning
            max_retries=2
        )
        
        # Create cleaning prompt
        cleaning_prompt = f"""Clean up this slide text that was extracted from a PDF. Remove:
- Math equation artifacts and garbled symbols
- Repeated headers/footers 
- Random formatting characters
- Broken table layouts
- Page numbers and slide metadata

Keep the meaningful content and structure. Make it readable and coherent:

---SLIDE TEXT---
{raw_content}
---END SLIDE TEXT---

Return only the cleaned text, no explanations."""

        logger.info(f"Cleaning {len(raw_content)} characters of slide content")
        
        # Get cleaned content from LLM
        cleaned_result = llm.invoke(cleaning_prompt)
        
        # Extract the actual text response
        if hasattr(cleaned_result, 'content'):
            cleaned_content = cleaned_result.content
        elif isinstance(cleaned_result, str):
            cleaned_content = cleaned_result
        else:
            cleaned_content = str(cleaned_result)
        
        cleaned_content = cleaned_content.strip()
        
        # Basic validation - make sure we didn't lose too much content
        if len(cleaned_content) < len(raw_content) * 0.3:  # Lost more than 70%
            logger.warning("LLM cleaning removed too much content, using original")
            return raw_content
        
        logger.info(f"Cleaned content: {len(raw_content)} → {len(cleaned_content)} characters")
        return cleaned_content
        
    except Exception as e:
        logger.error(f"Error cleaning slide content: {e}")
        logger.info("Falling back to original content")
        return raw_content


def clean_slides_text(slides: list) -> list:
    """
    Clean content for multiple slides while preserving structure
    
    Args:
        slides: List of slide dicts with 'content' field
    
    Returns:
        List of slides with cleaned content
    """
    cleaned_slides = []
    
    for slide in slides:
        if isinstance(slide, dict) and 'content' in slide:
            cleaned_slide = slide.copy()
            cleaned_slide['content'] = clean_slide_text(slide['content'])
            cleaned_slides.append(cleaned_slide)
        else:
            cleaned_slides.append(slide)  # Pass through unchanged
    
    return cleaned_slides


def get_cleaned_text(content_dict: dict) -> dict:
    """
    Clean both individual slides and full content in a slide extraction result
    
    Args:
        content_dict: Result from get_slide_text()
    
    Returns:
        Same dict with cleaned content and markdown format
    """
    if not content_dict.get('success'):
        return content_dict
    
    result = content_dict.copy()
    
    # Clean full content
    if 'content' in result:
        result['content'] = clean_slide_text(result['content'])
    
    # Clean individual slides
    if 'slides' in result:
        result['slides'] = clean_slides_text(result['slides'])
    
    # Create markdown version with hierarchy
    if 'slides' in result:
        images_data = result.get('images', [])
        result['content_markdown'] = format_slides_as_markdown(result['slides'], images_data)
    
    # Extract resources from content
    resources = extract_resources_from_content(result.get('content', ''))
    
    # Add QR code URLs if visual analysis was performed and QR codes found
    if 'images' in result and result.get('images'):
        try:
            from far_comms.utils.visual_analyzer import detect_qr_codes_in_images
            qr_codes = detect_qr_codes_in_images(result['images'])
            
            # Add QR URLs to resources
            for qr in qr_codes:
                if qr.get('type') == 'url':
                    url = qr.get('url')
                    page_num = qr.get('page', '?')
                    resources.append({
                        'name': f"QR Code (Slide {page_num})",
                        'url': url,
                        'context': f"QR code found on slide {page_num}"
                    })
                    
        except ImportError:
            pass  # QR detection not available
        except Exception as e:
            logger.warning(f"Error detecting QR codes for resources: {e}")
    
    if resources:
        result['resources'] = resources
        result['resources_formatted'] = format_resources_for_coda(resources)
    
    logger.info(f"Cleaned slide extraction for {result.get('file_name', 'unknown file')}")
    
    return result


def format_slides_as_markdown(slides: list, images_data: list = None) -> str:
    """
    Convert cleaned slides into structured markdown with hierarchy
    
    Args:
        slides: List of cleaned slide dicts with 'content' field
        images_data: Optional list of image dicts for adding image placeholders
    
    Returns:
        Markdown formatted string with slide hierarchy and image indicators
    """
    if not slides:
        return ""
    
    markdown_parts = []
    
    for slide in slides:
        if isinstance(slide, dict):
            page_num = slide.get('page', slide.get('slide', '?'))
            content = slide.get('content', '').strip()
            
            if content and content != f"[PAGE {page_num} - NO TEXT EXTRACTED]":
                # Add horizontal rule as page break (except for first slide)
                if markdown_parts:  # Not the first slide
                    markdown_parts.append("---")
                    markdown_parts.append("")  # Empty line after HR
                
                # Add image placeholder if images are available for this slide
                if images_data:
                    image_placeholder = get_image_placeholder(page_num, content, images_data)
                    if image_placeholder:
                        markdown_parts.append(image_placeholder)
                        markdown_parts.append("")  # Empty line after image
                
                # Process content to preserve structure
                formatted_content = format_slide_content(content)
                markdown_parts.append(formatted_content)
                markdown_parts.append("")  # Empty line after content
    
    return "\n".join(markdown_parts).strip()


def format_slide_content(content: str) -> str:
    """
    Format individual slide content with basic markdown structure
    
    Args:
        content: Cleaned slide content string
    
    Returns:
        Content formatted with basic markdown structure
    """
    if not content:
        return ""
    
    lines = content.split('\n')
    formatted_lines = []
    
    for line in lines:
        line = line.strip()
        if not line:
            formatted_lines.append("")
            continue
        
        # Convert bullet points and dashes to markdown bullets
        if line.startswith(('•', '▪', '▫', '‣', '⁃')) or line.startswith('- '):
            formatted_lines.append(f"- {line[1:].strip()}")
        elif line.startswith(('*', '·')) and len(line) > 2:
            formatted_lines.append(f"- {line[1:].strip()}")
        # Convert numbered lists
        elif line and line[0].isdigit() and len(line) > 2 and line[1:3] in ['. ', ') ']:
            formatted_lines.append(line)  # Keep numbered lists as-is
        # Potential headers (short lines, often all caps or title case)
        elif len(line) < 80 and (line.isupper() or line.istitle()) and not line.endswith('.'):
            formatted_lines.append(f"### {line}")
        else:
            formatted_lines.append(line)
    
    return "\n".join(formatted_lines)


def get_image_placeholder(page_num: int, content: str, images_data: list) -> str:
    """
    Generate image placeholder text based on slide content and available images
    
    Args:
        page_num: Slide/page number (1-indexed)
        content: Cleaned text content from slide
        images_data: List of image dicts with metadata
    
    Returns:
        Image placeholder string or empty string if no image
    """
    # Find corresponding image data
    image_info = None
    for img in images_data:
        if img.get('page') == page_num:
            image_info = img
            break
    
    if not image_info:
        return ""
    
    # Generate brief description based on content analysis
    description = infer_image_content(content)
    
    return f"*[Image: {description}]*"


def infer_image_content(content: str) -> str:
    """
    Infer what type of visual content might be present based on text content
    
    Args:
        content: Cleaned slide text content
    
    Returns:
        Brief description of likely image content
    """
    if not content:
        return "visual content"
    
    content_lower = content.lower()
    
    # Chart/graph indicators
    chart_keywords = ['chart', 'graph', 'plot', 'figure', 'data', 'results', 'analysis', 'percentage', '%', 'trend']
    if any(keyword in content_lower for keyword in chart_keywords):
        if 'result' in content_lower or 'analysis' in content_lower:
            return "chart showing results/analysis"
        elif any(word in content_lower for word in ['trend', 'over time', 'growth']):
            return "trend chart/graph"
        else:
            return "chart/graph with data"
    
    # Diagram/architecture indicators
    diagram_keywords = ['architecture', 'system', 'model', 'framework', 'process', 'workflow', 'pipeline']
    if any(keyword in content_lower for keyword in diagram_keywords):
        return "system/process diagram"
    
    # Screenshot/interface indicators  
    ui_keywords = ['interface', 'screen', 'demo', 'example', 'tool', 'platform']
    if any(keyword in content_lower for keyword in ui_keywords):
        return "interface/demo screenshot"
    
    # Mathematical/technical indicators
    math_keywords = ['equation', 'formula', 'algorithm', 'method', 'approach', 'technique']
    if any(keyword in content_lower for keyword in math_keywords):
        return "technical diagram/equations"
    
    # Default fallback
    return "slide visual content"


def extract_resources_from_content(content: str) -> list:
    """
    Extract references to papers, websites, and resources from slide content
    
    Args:
        content: Full slide content text
    
    Returns:
        List of resource dicts with 'name' and 'url' keys
    """
    if not content:
        return []
    
    resources = []
    lines = content.split('\n')
    
    # URL patterns
    import re
    url_pattern = r'https?://[^\s\)]+|www\.[^\s\)]+|\b[a-zA-Z0-9.-]+\.(com|org|edu|net|gov|io|ai|co\.uk|ac\.uk)\b[^\s\)]*'
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
        
        # Find URLs in the line
        urls = re.findall(url_pattern, line, re.IGNORECASE)
        
        for url in urls:
            # Clean up URL (remove trailing punctuation)
            clean_url = url.rstrip('.,;:!?)')
            
            # Skip if URL is too short or looks malformed
            if len(clean_url) < 8:
                continue
            
            # Add protocol if missing
            if not clean_url.startswith(('http://', 'https://')):
                if clean_url.startswith('www.'):
                    clean_url = 'https://' + clean_url
                else:
                    clean_url = 'https://' + clean_url
            
            # Extract resource name from context
            resource_name = extract_resource_name(line, url)
            
            resources.append({
                'name': resource_name,
                'url': clean_url,
                'context': line[:100]  # Keep some context
            })
    
    # Remove duplicates based on URL
    seen_urls = set()
    unique_resources = []
    for resource in resources:
        if resource['url'] not in seen_urls:
            seen_urls.add(resource['url'])
            unique_resources.append(resource)
    
    return unique_resources


def extract_resource_name(line: str, url: str) -> str:
    """
    Extract a meaningful name for a resource from the line context
    
    Args:
        line: The line containing the URL
        url: The URL found in the line
    
    Returns:
        Descriptive name for the resource
    """
    # Remove the URL from the line to get potential name
    line_without_url = line.replace(url, '').strip()
    
    # Clean up common prefixes/suffixes
    line_clean = re.sub(r'^[-•*\s]+|[-•*\s]+$', '', line_without_url)
    line_clean = re.sub(r'^(source:|reference:|see:|from:|paper:|study:|website:|link:)\s*', '', line_clean, flags=re.IGNORECASE)
    
    # If we have meaningful text, use it
    if line_clean and len(line_clean.strip()) > 3:
        # Limit length and clean up
        name = line_clean.strip()[:60]
        if len(line_clean) > 60:
            name = name.rsplit(' ', 1)[0] + '...'  # Break at word boundary
        return name
    
    # Extract domain name as fallback
    try:
        from urllib.parse import urlparse
        parsed = urlparse(url if url.startswith(('http://', 'https://')) else 'https://' + url)
        domain = parsed.netloc.lower()
        
        # Clean up common prefixes
        domain = re.sub(r'^www\.', '', domain)
        
        # Capitalize for readability
        return domain.replace('.', ' ').title().replace(' Com', '.com').replace(' Org', '.org')
        
    except:
        return "Resource"


def format_resources_for_coda(resources: list) -> str:
    """
    Format resources for Coda "Resources" column in specified format
    
    Args:
        resources: List of resource dicts with 'name' and 'url'
    
    Returns:
        Formatted string: "{resource_name} - {resource_url}"
    """
    if not resources:
        return ""
    
    formatted_items = []
    for resource in resources:
        name = resource.get('name', 'Resource')
        url = resource.get('url', '')
        
        if url:  # Only include if we have a valid URL
            formatted_items.append(f"{name} - {url}")
    
    return '\n'.join(formatted_items)


# Example usage and testing
if __name__ == "__main__":
    # Test with the provided file
    test_file = "/Users/cheng2/Desktop/agents/far_comms/data/slides/11_50_Xiaoyuan Yi-ValueCompass_updated.pptx.pdf"
    
    print("Testing slide extraction...")
    result = get_slide_text(test_file)
    
    print(f"Success: {result.get('success')}")
    print(format_extraction_summary(result))
    
    if result.get("success"):
        print(f"\nFirst 500 characters of content:")
        print(result.get("content", "")[:500])
        print("\n" + "="*50)
        
        print(f"\nFirst slide/page content:")
        slides = result.get("slides", [])
        if slides:
            print(slides[0].get("content", "")[:300])
    else:
        print(f"Error: {result.get('error')}")